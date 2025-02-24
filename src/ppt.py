from pptx import Presentation
from pptx.util import Inches
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class GeneratePPT:
    @staticmethod
    def create_powerpoint(screenshot_paths, template_path, output_path):
        try:
            # Load the template
            prs = Presentation(template_path)
            
            # Skip first slide (intro) and add screenshots to slides 2 and 3
            for i, screenshot_path in enumerate(screenshot_paths):
                if i < 2:  # Only process first two screenshots
                    slide = prs.slides[i + 1]  # +1 to skip intro slide
                    
                    # Add screenshot to the empty slide
                    # Center the screenshot on the slide
                    left = Inches(1)  # 1 inch margin from left
                    top = Inches(1)   # 1 inch margin from top
                    width = Inches(8)  # 8 inches wide
                    height = Inches(5.5)  # 5.5 inches high
                    
                    slide.shapes.add_picture(
                        screenshot_path,
                        left,
                        top,
                        width,
                        height
                    )
                    logger.info(f"Added screenshot to slide {i + 2}")
            
            # Save the presentation
            prs.save(output_path)
            logger.info(f"PowerPoint saved: {output_path}")
            return True
        except Exception as e:
            logger.error(f"PowerPoint creation failed: {e}")
            return False